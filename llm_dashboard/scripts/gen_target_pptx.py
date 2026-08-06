# -*- coding: utf-8 -*-
"""
从 dashboard/data/tg_data.js 解析目标达成数据，筛出 2026-07，生成 PPTX 汇报。
用法: python scripts/gen_target_pptx.py
"""
import json, os, re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_TICK_LABEL_POSITION

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC = os.path.join(ROOT, 'dashboard', 'data', 'tg_data.js')
OUT_DIR = os.path.join(ROOT, 'exports')
os.makedirs(OUT_DIR, exist_ok=True)
OUT = os.path.join(OUT_DIR, '目标达成看板_2026年7月_含负责人国家费比.pptx')

# ---------- 配色 ----------
BLUE   = RGBColor(0x33, 0x60, 0x8c)
ORANGE = RGBColor(0xd9, 0x9a, 0x3d)
GREEN  = RGBColor(0x3a, 0x8f, 0x6c)
RED    = RGBColor(0xc0, 0x45, 0x3e)
INK    = RGBColor(0x26, 0x2b, 0x33)
GREY   = RGBColor(0x7a, 0x80, 0x89)
LIGHT  = RGBColor(0xf2, 0xf4, 0xf7)
WHITE  = RGBColor(0xff, 0xff, 0xff)

# ---------- 工具 ----------
def money(v):
    if v is None or (isinstance(v, float) and (v != v)):  # NaN
        return '—'
    return '$' + format(v, ',.0f')

def pct(v, d=1):
    if v is None or (isinstance(v, float) and v != v):
        return '—'
    return f"{v*100:.{d}f}%"

def num(v):
    if v is None or (isinstance(v, float) and v != v):
        return '—'
    return format(v, ',.0f')

def load_data():
    txt = open(SRC, encoding='utf-8').read()
    i = txt.index('{', txt.index('window.TG_DATA'))
    depth = 0
    for j in range(i, len(txt)):
        if txt[j] == '{': depth += 1
        elif txt[j] == '}':
            depth -= 1
            if depth == 0:
                end = j + 1; break
    return json.loads(txt[i:end])

def month_rows(daily, DIM, y, mo):
    """返回某年某月的 daily 行（按日期前缀），以及该月日期列表"""
    out = []
    for r in daily:
        ds = DIM['d'][r[0]]
        if ds.startswith(f"{y:04d}-{mo:02d}"):
            out.append(r)
    return out

def sum_d(rows):
    t = dict(sales=0.0, adsp=0.0, adsa=0.0, imp=0.0, clk=0.0, adod=0.0, od=0.0)
    for r in rows:
        t['sales'] += r[5]; t['adsp'] += r[6]; t['adsa'] += r[7]
        t['imp']   += r[8]; t['clk']  += r[9]; t['adod'] += r[10]; t['od'] += r[11]
    return t

def sum_t(rows):
    ts = tp = 0.0
    for r in rows:
        ts += r[6]; tp += r[7]
    return ts, tp

def metrics_d(t):
    adsp, sales, adsa, clk, imp, adod = t['adsp'], t['sales'], t['adsa'], t['clk'], t['imp'], t['adod']
    return {
        '广告花费': adsp,
        '总销售额': sales,
        '广告销售额': adsa,
        'TACOS': (adsp/sales) if sales > 0 else float('nan'),
        'ACOS': (adsp/adsa) if adsa > 0 else float('nan'),
        'CPC': (adsp/clk) if clk > 0 else float('nan'),
        'CTR': (clk/imp) if imp > 0 else float('nan'),
        'CVR': (adod/clk) if clk > 0 else float('nan'),
        '点击量': clk,
        '订单量': adod,
    }

def by_dim(rows, targets, DIM, dim_idx, tgt_dim_idx, name_key, y, mo):
    """按维度(owner/shop/country/cat)聚合 实际 与 目标，返回 list of dict"""
    act_map, tgt_map = {}, {}
    name_arr = DIM[name_key]
    for r in rows:
        k = r[dim_idx]
        a = act_map.get(k) or {'sales':0.0,'adsp':0.0,'adsa':0.0}
        a['sales'] += r[5]; a['adsp'] += r[6]; a['adsa'] += r[7]
        act_map[k] = a
    for r in targets:
        if r[0] != y or r[1] != mo: continue
        k = r[tgt_dim_idx]
        b = tgt_map.get(k) or {'ts':0.0,'tp':0.0}
        b['ts'] += r[6]; b['tp'] += r[7]
        tgt_map[k] = b
    res = []
    keys = set(act_map) | set(tgt_map)
    for k in keys:
        a = act_map.get(k, {'sales':0,'adsp':0,'adsa':0})
        b = tgt_map.get(k, {'ts':0,'tp':0})
        tt = (b['tp']/b['ts']) if b['ts'] > 0 else float('nan')
        at = (a['adsp']/a['sales']) if a['sales'] > 0 else float('nan')
        spend = (a['adsp']/b['tp']) if b['tp'] > 0 else float('nan')
        tacos = (tt/at) if (at and at > 0 and tt == tt) else float('nan')
        res.append({'name': name_arr[k] if isinstance(k,int) else k, 'actSp': a['adsp'], 'tgtSp': b['tp'],
                    'actSales': a['sales'], 'spend': spend, 'tt': tt, 'at': at, 'tacos': tacos})
    return res

# ---------- 加载 ----------
DATA = load_data()
DIM = DATA['dims']
DAILY = DATA['daily']
TGT = DATA['targets']

Y, MO = 2026, 7
july_rows = month_rows(DAILY, DIM, Y, MO)
june_rows = month_rows(DAILY, DIM, 2026, 6)

j_act = sum_d(july_rows)
j_ts, j_tp = sum_t([r for r in TGT if r[0]==Y and r[1]==MO])
june_act = sum_d(june_rows)
june_ts, june_tp = sum_t([r for r in TGT if r[0]==2026 and r[1]==6])

j_spend = (j_act['adsp']/j_tp) if j_tp > 0 else float('nan')
j_tgtTacos = (j_tp/j_ts) if j_ts > 0 else float('nan')
j_actTacos = (j_act['adsp']/j_act['sales']) if j_act['sales'] > 0 else float('nan')
j_tacos = (j_tgtTacos/j_actTacos) if (j_actTacos and j_actTacos>0 and j_tgtTacos==j_tgtTacos) else float('nan')
j_time = 1.0  # 7月整月已结束
j_dev = (j_spend - j_time) if j_spend==j_spend else float('nan')

# 维度聚合
cats = by_dim(july_rows, TGT, DIM, 4, 5, 'g', Y, MO)
owners = by_dim(july_rows, TGT, DIM, 1, 2, 'o', Y, MO)
countries = by_dim(july_rows, TGT, DIM, 3, 4, 'c', Y, MO)

# 每日序列（7月）
july_dates = [d for d in DIM['d'] if d.startswith('2026-07')]
date_idx = {d: i for i, d in enumerate(DIM['d'])}
daily_sp, daily_tc = [], []
for d in july_dates:
    rows = [r for r in july_rows if DIM['d'][r[0]] == d]
    t = sum_d(rows)
    daily_sp.append(t['adsp'])
    daily_tc.append((t['adsp']/t['sales']) if t['sales']>0 else 0.0)

# ---------- 打印校验 ----------
print("=== 7月总览 ===")
print("目标总销售额", money(j_ts), "目标广告花费", money(j_tp))
print("实际总销售额", money(j_act['sales']), "实际广告花费", money(j_act['adsp']), "实际广告销售额", money(j_act['adsa']))
print("花费完成率", pct(j_spend), "费比完成率", pct(j_tacos), "目标TACOS", pct(j_tgtTacos), "实际TACOS", pct(j_actTacos))
print("ACOS", pct(j_act['adsp']/j_act['adsa'] if j_act['adsa']>0 else float('nan')), "CPC", f"${j_act['adsp']/j_act['clk']:.2f}" if j_act['clk']>0 else '—')
print("类目数", len(cats), "负责人数", len(owners), "国家数", len(countries))
print("类目达标(top5):", [(c['name'], round(c['tacos']*100,1)) for c in sorted(cats, key=lambda x:-x['tacos'])[:5]])

# ================= PPTX 构建 =================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def add_rect(slide, x, y, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(1)
    sp.shadow.inherit = False
    return sp

def set_text(tf, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = font
    # 设置东亚字体
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}ea')
    if ea is None:
        ea = rPr.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}ea', {})
        rPr.append(ea)
    ea.set('typeface', font)
    return p

def add_textbox(slide, x, y, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    set_text(tf, text, size, color, bold, align)
    return tb

def title_bar(slide, text, sub=None):
    bar = add_rect(slide, 0, 0, SW, Inches(1.0), BLUE)
    add_textbox(slide, Inches(0.5), Inches(0.12), Inches(11), Inches(0.5), text, 26, WHITE, True)
    if sub:
        add_textbox(slide, Inches(0.5), Inches(0.6), Inches(12), Inches(0.35), sub, 13, RGBColor(0xcf,0xdd,0xea))

# ---- Slide 1: 封面 ----
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SW, SH, BLUE)
add_rect(s, 0, Inches(4.6), SW, Inches(0.08), ORANGE)
add_textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.0), "跨境电商广告 · 目标达成看板", 40, WHITE, True)
add_textbox(s, Inches(0.9), Inches(3.1), Inches(11.5), Inches(0.7), "2026 年 7 月达成数据汇报", 26, RGBColor(0xe6,0xee,0xf5))
add_textbox(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(0.5),
            f"数据区间：2026-07-01 ~ 2026-07-31    |    生成于本地看板数据 v20260803d1", 14, RGBColor(0xb9,0xc9,0xd9))
add_textbox(s, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.5),
            f"实际广告花费 {money(j_act['adsp'])}  /  目标 {money(j_tp)}  /  花费完成率 {pct(j_spend)}", 15, RGBColor(0xe6,0xee,0xf5), True)

# ---- Slide 2: 总览 KPI ----
s = prs.slides.add_slide(BLANK)
title_bar(s, "一、7 月目标达成总览", "核心指标：花费完成率、费比（TACOS）完成率、时间进度偏差")
ac = j_act
kpis = [
    ("实际广告花费", money(ac['adsp']), f"目标 {money(j_tp)}", LIGHT, BLUE),
    ("花费完成率", pct(j_spend), "vs 时间进度 100%", LIGHT, GREEN if (j_spend==j_spend and abs(j_spend-1)<=0.05) else (GREEN if j_spend>=1 else RED)),
    ("实际 TACOS", pct(j_actTacos), f"目标 {pct(j_tgtTacos)}", LIGHT, GREEN if (j_actTacos<=j_tgtTacos) else RED),
    ("费比完成率", pct(j_tacos), "≥100% 达标", LIGHT, GREEN if (j_tacos==j_tacos and j_tacos>=1) else RED),
    ("实际 ACOS", pct(ac['adsp']/ac['adsa'] if ac['adsa']>0 else float('nan')), f"广告销售额 {money(ac['adsa'])}", LIGHT, BLUE),
    ("实际 CPC", f"${ac['adsp']/ac['clk']:.2f}" if ac['clk']>0 else '—', f"点击 {num(ac['clk'])}", LIGHT, BLUE),
    ("总销售额", money(ac['sales']), f"目标 {money(j_ts)}", LIGHT, BLUE),
    ("进度偏差", ("+" if j_dev>0 else "") + pct(j_dev) if j_dev==j_dev else '—', "花费 vs 时间", LIGHT,
        GREEN if (j_dev==j_dev and abs(j_dev)<=0.05) else (ORANGE if j_dev<0 else RED)),
]
gx, gy = Inches(0.5), Inches(1.25)
cw, ch = Inches(3.0), Inches(1.7)
gx_step, gy_step = Inches(3.18), Inches(1.85)
for idx, (lab, val, sub, bg, fg) in enumerate(kpis):
    r, c = divmod(idx, 4)
    x = gx + c*gx_step; y = gy + r*gy_step
    card = add_rect(s, x, y, cw, ch, bg)
    set_text(card.text_frame, "", 1)  # init
    tf = card.text_frame; tf.word_wrap = True
    p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run(); r1.text = lab; r1.font.size = Pt(13); r1.font.color.rgb = GREY; r1.font.name='Microsoft YaHei'
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = val; r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = fg; r2.font.name='Microsoft YaHei'
    p3 = tf.add_paragraph(); p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run(); r3.text = sub; r3.font.size = Pt(11); r3.font.color.rgb = GREY; r3.font.name='Microsoft YaHei'

# ---- Slide 3: 类目费比完成率 ----
s = prs.slides.add_slide(BLANK)
title_bar(s, "二、类目费比（TACOS）完成率", "完成率 = 目标 TACOS / 实际 TACOS，≥100% 表示费比达标（实际费比不高于目标）")
cats_sorted = sorted([c for c in cats if c['tacos']==c['tacos']], key=lambda x:-x['tacos'])
top = cats_sorted[:12]
# 图表
chart_data = CategoryChartData()
chart_data.categories = [c['name'] for c in top][::-1]
chart_data.add_series('费比完成率', tuple(round(c['tacos']*100,1) for c in top)[::-1])
gf = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.5), Inches(1.2), Inches(7.6), Inches(5.9), chart_data)
ch = gf.chart
ch.has_legend = False
ch.has_title = True; ch.chart_title.text_frame.text = "各类目费比完成率（%）"
ch.chart_title.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
plot = ch.plots[0]
plot.has_data_labels = True
plot.data_labels.number_format = '0.0"%"'; plot.data_labels.number_format_is_linked = False
plot.data_labels.font.size = Pt(9)
ser = plot.series[0]
ser.format.fill.solid(); ser.format.fill.fore_color.rgb = BLUE
# 右表
tx = Inches(8.4); tw = Inches(4.4)
add_textbox(s, tx, Inches(1.2), tw, Inches(0.4), "类目排行（完成率降序）", 14, INK, True)
tbl_shape = s.shapes.add_table(len(top[:10])+1, 3, tx, Inches(1.65), tw, Inches(5.0))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.0); tbl.columns[1].width = Inches(1.2); tbl.columns[2].width = Inches(1.2)
hdr = ['类目','目标TACOS','实际TACOS']
for j,h in enumerate(hdr):
    cell = tbl.cell(0,j); cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
    p = cell.text_frame.paragraphs[0]; p.runs[0].font.size=Pt(11); p.runs[0].font.bold=True; p.runs[0].font.color.rgb=WHITE
for i, c in enumerate(top[:10], start=1):
    vals = [c['name'], pct(c['tt']), pct(c['at'])]
    for j,v in enumerate(vals):
        cell = tbl.cell(i,j); cell.text = str(v)
        p = cell.text_frame.paragraphs[0]; p.runs[0].font.size=Pt(10)
        if j==0: p.runs[0].font.color.rgb=INK
        else: p.runs[0].font.color.rgb=(GREEN if (j==2 and c['at']<=c['tt']) else (RED if (j==2 and c['at']>c['tt']) else GREY))

# ---- Slide 4: 负责人 & 国家 费比完成率 ----
s = prs.slides.add_slide(BLANK)
title_bar(s, "三、负责人 / 国家 费比（TACOS）完成率排行", "费比完成率 = 目标 TACOS / 实际 TACOS，≥100% 表示费比达标（实际费比不高于目标）；末列附花费完成率对照")
def tacos_table(slide, x, y, w, rows, dim_name):
    rows = sorted([r for r in rows if r['tacos']==r['tacos']], key=lambda z:-z['tacos'])
    rows = rows[:8]
    add_textbox(slide, x, y, w, Inches(0.4), f"{dim_name} Top（按费比完成率降序）", 14, INK, True)
    t = slide.shapes.add_table(len(rows)+1, 5, x, y+Inches(0.45), w, Inches(4.4)).table
    t.columns[0].width = Inches(1.5); t.columns[1].width = Inches(1.0); t.columns[2].width = Inches(1.0); t.columns[3].width = Inches(1.4); t.columns[4].width = Inches(1.1)
    for j,h in enumerate(['维度','目标TACOS','实际TACOS','费比完成率','花费完成率']):
        cell=t.cell(0,j); cell.text=h; cell.fill.solid(); cell.fill.fore_color.rgb=BLUE
        r0=cell.text_frame.paragraphs[0].runs[0]; r0.font.size=Pt(10); r0.font.bold=True; r0.font.color.rgb=WHITE; r0.font.name='Microsoft YaHei'
    for i,rr in enumerate(rows, start=1):
        vals=[rr['name'], pct(rr['tt']), pct(rr['at']), pct(rr['tacos']), pct(rr['spend'])]
        for j,v in enumerate(vals):
            cell=t.cell(i,j); cell.text=str(v)
            p=cell.text_frame.paragraphs[0]; p.runs[0].font.size=Pt(10); p.runs[0].font.name='Microsoft YaHei'
            if j==0:
                p.runs[0].font.color.rgb=INK
            elif j==3:
                sp=rr['tacos']
                p.runs[0].font.color.rgb = GREEN if sp>=1 else (ORANGE if sp>=0.85 else RED)
                p.runs[0].font.bold=True
            elif j==4:
                sp=rr['spend']
                p.runs[0].font.color.rgb = GREEN if sp>=1 else (ORANGE if sp>=0.85 else RED)
                p.runs[0].font.bold=True
            else:
                p.runs[0].font.color.rgb=GREY
tacos_table(s, Inches(0.5), Inches(1.2), Inches(6.0), owners, "负责人")
tacos_table(s, Inches(6.9), Inches(1.2), Inches(6.0), countries, "国家")

# ---- Slide 5: 每日趋势 ----
s = prs.slides.add_slide(BLANK)
title_bar(s, "四、7 月每日广告花费与 TACOS 趋势", "观察花费投放节奏与费比波动")
# 左：每日花费
cd1 = CategoryChartData()
cd1.categories = [d[5:] for d in july_dates]
cd1.add_series('每日广告花费', tuple(round(x,1) for x in daily_sp))
gf1 = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.2), Inches(6.1), Inches(5.6), cd1)
c1 = gf1.chart; c1.has_legend=False; c1.has_title=True
c1.chart_title.text_frame.text="每日广告花费 ($)"
c1.chart_title.text_frame.paragraphs[0].runs[0].font.size=Pt(12)
c1.plots[0].series[0].format.fill.solid(); c1.plots[0].series[0].format.fill.fore_color.rgb=BLUE
# 右：每日 TACOS
cd2 = CategoryChartData()
cd2.categories = [d[5:] for d in july_dates]
cd2.add_series('每日 TACOS', tuple(round(x*100,2) for x in daily_tc))
gf2 = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(6.9), Inches(1.2), Inches(6.0), Inches(5.6), cd2)
c2 = gf2.chart; c2.has_legend=False; c2.has_title=True
c2.chart_title.text_frame.text="每日 TACOS (%)"
c2.chart_title.text_frame.paragraphs[0].runs[0].font.size=Pt(12)
c2.plots[0].series[0].format.line.color.rgb=ORANGE; c2.plots[0].series[0].format.line.width=Pt(2)
c2.plots[0].has_data_labels=False

# ---- Slide 6: 差距归因 (7月 vs 6月) ----
s = prs.slides.add_slide(BLANK)
title_bar(s, "五、差距归因：7 月 vs 6 月", "对比上月同口径，定位花费/费比变化驱动因素")
j_m = metrics_d(j_act); n_m = metrics_d(june_act)
order = ['广告花费','总销售额','广告销售额','TACOS','ACOS','CPC','CTR','CVR','点击量']
def fmt(name, v):
    if name in ('TACOS','ACOS','CTR','CVR'): return pct(v)
    if name=='CPC': return f"${v:.2f}" if v==v else '—'
    return money(v)
rows_tbl = []
for name in order:
    jv, nv = j_m[name], n_m[name]
    if (jv==jv) and (nv==nv) and nv not in (0,):
        if name in ('TACOS','ACOS','CTR','CVR'):
            d = (jv - nv)*100
            dtxt = ("+" if d>0 else "") + f"{d:.1f}pp"
        else:
            d = (jv/nv - 1)*100
            dtxt = ("+" if d>0 else "") + f"{d:.1f}%"
    else:
        dtxt = '—'
    rows_tbl.append((name, fmt(name,jv), fmt(name,nv), dtxt))
t = s.shapes.add_table(len(rows_tbl)+1, 4, Inches(0.5), Inches(1.25), Inches(12.3), Inches(5.3)).table
t.columns[0].width=Inches(3.3); t.columns[1].width=Inches(3.0); t.columns[2].width=Inches(3.0); t.columns[3].width=Inches(3.0)
for j,h in enumerate(['指标','7月','6月','变化']):
    cell=t.cell(0,j); cell.text=h; cell.fill.solid(); cell.fill.fore_color.rgb=BLUE
    r0=cell.text_frame.paragraphs[0].runs[0]; r0.font.size=Pt(12); r0.font.bold=True; r0.font.color.rgb=WHITE
for i,(nm,jv,nv,d) in enumerate(rows_tbl, start=1):
    for j,v in enumerate([nm,jv,nv,d]):
        cell=t.cell(i,j); cell.text=str(v)
        p=cell.text_frame.paragraphs[0]; p.runs[0].font.size=Pt(11)
        if j==0: p.runs[0].font.color.rgb=INK; p.runs[0].font.bold=True
        elif j==3:
            if d not in ('—',''):
                good = (d.startswith('+') and nm in ('总销售额','广告销售额','CTR','CVR','点击量')) or \
                       (d.startswith('-') and nm in ('TACOS','ACOS','CPC'))
                p.runs[0].font.color.rgb = GREEN if good else RED
                p.runs[0].font.bold=True

# ---- Slide 7: 改进建议 ----
s = prs.slides.add_slide(BLANK)
title_bar(s, "六、改进建议", "基于 7 月达成与上月对比自动生成")
advice = []
if j_dev == j_dev:
    if j_dev < -0.05:
        advice.append(("⏩ 花费投放提速", f"花费完成率 {pct(j_spend)} 落后时间进度 100% 约 {pct(abs(j_dev))}。建议：① 核查每日预算是否提前触顶，触顶活动加预算 20%~30%；② 对优质词提高竞价抢量；③ 补充新词/新活动扩大覆盖。"))
    elif j_dev > 0.05:
        over = max(j_act['adsp']/max(j_time,0.01) - j_tp, 0)
        advice.append(("🛑 控制投放节奏", f"花费完成率超前，按当前节奏月底将超支约 {money(over)}。建议对高 ACOS 活动下调预算与竞价，回收低效花费。"))
    else:
        advice.append(("✅ 投放节奏健康", f"花费完成率与时间进度偏差仅 {pct(j_dev)}，保持当前节奏，重点转向结构性优化（词级调价、时段分配）。"))
if j_tacos == j_tacos:
    if j_tacos < 0.85:
        advice.append(("🚨 费比明显超标，优先控费", f"实际 TACOS {pct(j_actTacos)} 显著高于目标 {pct(j_tgtTacos)}（完成率 {pct(j_tacos)}）。行动：① 暂停零转化高花费词；② 高 ACOS 词按幅度降竞价；③ 排查低自然流量 SKU 的 Listing 权重。"))
    elif j_tacos < 1:
        advice.append(("⚠️ 费比小幅超标，结构调优", f"实际 TACOS {pct(j_actTacos)} 略高于目标 {pct(j_tgtTacos)}。差距不大，压缩广泛匹配低效流量、优化否定词库、预算向达标类目/国家倾斜即可收敛。"))
    else:
        advice.append(("🏆 费比达标，可适度放量", f"实际 TACOS {pct(j_actTacos)} 优于目标 {pct(j_tgtTacos)}，有 {pct(max(j_tacos-1,0))} 余量。在守住费比红线的前提下，对优质词与达标类目增预算换规模。"))
# 类目洞察
if cats_sorted:
    best = cats_sorted[0]; worst = cats_sorted[-1]
    advice.append(("📊 类目结构差异", f"费比最优类目「{best['name']}」完成率 {pct(best['tacos'])}；最弱类目「{worst['name']}」完成率 {pct(worst['tacos'])}。建议预算向达标类目集中，最弱类目先做词级诊断再决定加投或收缩。"))
advice.append(("📅 月末冲刺机制", "每周一复盘各负责人/国家花费与费比完成率排行，落后者提交改进计划；月末最后一周锁定预算分配，避免为冲花费完成率而牺牲费比。"))

y = Inches(1.25)
for ico_tit, txt in advice:
    card = add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.92), LIGHT)
    tf = card.text_frame; tf.word_wrap=True; tf.margin_left=Inches(0.15); tf.margin_top=Inches(0.05)
    p = tf.paragraphs[0]
    r=p.add_run(); r.text=ico_tit; r.font.size=Pt(13); r.font.bold=True; r.font.color.rgb=BLUE; r.font.name='Microsoft YaHei'
    p2=tf.add_paragraph(); r2=p2.add_run(); r2.text=txt; r2.font.size=Pt(11); r2.font.color.rgb=INK; r2.font.name='Microsoft YaHei'
    y = y + Inches(1.0)

prs.save(OUT)
print("\nSaved ->", OUT)
print("Slides:", len(prs.slides._sldIdLst))
